# watermarking all images with nike logo using python-imagemagick-wand
    # 1. download imagemagick from www.imagemagick.org 
    # 2. install wand - pip install wand
    # 3. importing library from Wand
from wand import image
from wand.image import Image
from wand.compat import nested
from wand.display import display

images = ["setup.jpg",
"flower pot.jpg",
"camera.jpg",
 "lemon tea.jpg ",
 'wood.jpg'
 ]

for item in images:
    logo = Image(filename="logo.png")
    image = Image(filename= item)
    image.composite_channel('all_channels', logo, 'dissolve',0,0)
    image.save(filename=item) 






# for install pptx ---- 
# pip install python-pptx

# for creating presentation ---- 
from pptx import Presentation
import pptx

# for inserting picture in slides ---
from pptx.util import Inches
import os
 
prs = Presentation()
 
class MySlide:
    def __init__(self, data):
        self.layout = prs.slide_layouts[data[3]]
        self.slide =  prs.slides.add_slide(self.layout)
        self.title =  self.slide.shapes.title
        self.title.text=data[0]
        self.subtitle=self.slide.placeholders[1]
        self.subtitle.text=data[1]
        if data[2] != "":
        	self.slide.shapes.add_picture(data[2],pptx.util.Inches(2), pptx.util.Inches(3),width=pptx.util.Inches(5), height=pptx.util.Inches(4))
 
slides = [
    ["Camera",      
     "Subtitle(1)",
     'camera.jpg',
     1],
    ["Swiss Cheese Plant",       
     "Subtitle(2)",
     "flower pot.jpg",
     1],
    ["Lemon Tea",       
     "Subtitle(3)",
     "lemon tea.jpg",
     1],
     ["Computer Setup",
     "Subtitle(4)",
     "setup.jpg",
     1],
     ["Brown Wooden Interior",
     "Subtitle(5)",
     "wood.jpg",
     1]
]
 
for each_slide in slides:
    MySlide(each_slide)
 
prs.save("presentation.pptx")
os.startfile("presentation.pptx")

